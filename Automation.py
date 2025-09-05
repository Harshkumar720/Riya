# Automation.py
import subprocess
import webbrowser
import os
import pyautogui
import time
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from groq import Groq
from PIL import Image
import speech_recognition as sr
import pyttsx3
import re  # ✅ added for normalization

# ------------------ Groq client setup ------------------
GROQ_API_KEY = ""
client = Groq(api_key=GROQ_API_KEY)

# ------------------ TTS Engine (Female Voice) ------------------
engine = pyttsx3.init()
voices = engine.getProperty('voices')
# Usually female voice is at index 1
engine.setProperty('voice', voices[1].id)
engine.setProperty('rate', 160)  # optional: adjust speed
engine.setProperty('volume', 1.0)

def speak(text):
    engine.say(text)
    engine.runAndWait()

# ------------------ Voice Recognition ------------------
recognizer = sr.Recognizer()
def listen_command():
    with sr.Microphone() as source:
        print("🎤 Listening for command...")
        recognizer.adjust_for_ambient_noise(source)
        audio = recognizer.listen(source, phrase_time_limit=6)
        try:
            command = recognizer.recognize_google(audio)
            print(f"🗣 You said: {command}")
            return command.lower()
        except:
            return ""

# ------------------ Groq Content Generator ------------------
def write_content(prompt, max_tokens=500):
    try:
        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": "You are an expert writer for applications, essays, reports, speeches, stories, and PPT slides."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=max_tokens,
            temperature=0.7
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"⚠️ Groq Error: {str(e)}"

# ------------------ Helper: Open Notepad ------------------
def open_notepad_with_content(content):
    temp_file = "temp_note.txt"
    with open(temp_file, "w", encoding="utf-8") as f:
        f.write(content)
    os.startfile(temp_file)

# ------------------ Web Image Fetch ------------------
def fetch_image_url(query):
    try:
        search_term = query.replace(" ", "+")
        return f"https://source.unsplash.com/800x600/?{search_term}"
    except:
        return None

# ------------------ PPT Creation ------------------
def create_ppt(topic, slides_count=10):
    try:
        prs = Presentation()
        for i in range(slides_count):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"{topic} - Slide {i+1}"

            slide_text = write_content(f"Write 4 concise bullet points for slide {i+1} of a presentation on {topic}.")
            slide.placeholders[1].text = slide_text

            img_url = fetch_image_url(f"{topic} slide {i+1}")
            if img_url:
                try:
                    img_data = requests.get(img_url).content
                    slide.shapes.add_picture(BytesIO(img_data), Inches(5), Inches(1.5), width=Inches(4))
                except:
                    pass

        filename = f"{topic.replace(' ', '_')}_{int(time.time())}.pptx"
        prs.save(filename)
        os.startfile(os.path.abspath(filename))
        return f"✅ PowerPoint created & opened: {filename}"
    except:
        return "⚠️ Failed to create PPT. Make sure python-pptx is installed."

# ------------------ App Open/Download ------------------
APP_STORE_LINKS = {
    "chrome": "https://www.google.com/chrome/",
    "firefox": "https://www.mozilla.org/firefox/new/",
    "edge": "https://www.microsoft.com/edge",
    "vlc": "https://www.videolan.org/vlc/",
    "spotify": "https://www.spotify.com/download/windows/",
    "word": "https://www.microsoft.com/microsoft-365/word",
    "excel": "https://www.microsoft.com/microsoft-365/excel",
    "powerpoint": "https://www.microsoft.com/microsoft-365/powerpoint",
    "notepad++": "https://notepad-plus-plus.org/downloads/",
    "zoom": "https://zoom.us/download",
    "teams": "https://www.microsoft.com/microsoft-teams/download-app",
    "jiohotstar": "https://www.hotstar.com/in",
    "youtube": "https://www.youtube.com/",
    "telegram": "https://desktop.telegram.org/",
    "whatsapp": "https://web.whatsapp.com/"
}

def find_installed_app(app_name):
    paths_to_check = [
        os.environ.get("ProgramFiles"),
        os.environ.get("ProgramFiles(x86)"),
        os.path.join(os.environ.get("LOCALAPPDATA", ""), "WhatsApp")
    ]
    for folder in paths_to_check:
        if folder:
            for root, dirs, files in os.walk(folder):
                for file in files:
                    if app_name.lower() in file.lower() and file.endswith(".exe"):
                        return os.path.join(root, file)
    return None

def open_application(app_name):
    app_path = find_installed_app(app_name)
    if app_path:
        try:
            subprocess.Popen(app_path)
            return f"✅ Opened application: {app_name}"
        except Exception as e:
            return f"⚠️ Could not open: {str(e)}"
    else:
        link = next((url for key, url in APP_STORE_LINKS.items() if key in app_name.lower()), None)
        if link:
            webbrowser.open(link)
            return f"{app_name} not installed. Opening official page..."
        else:
            webbrowser.open(f"https://www.microsoft.com/store/search/apps?q={app_name}")
            return f"{app_name} not installed. Searching in Microsoft Store..."

# ------------------ Close Applications (robust) ------------------

def _normalize_app_name(s: str) -> str:
    """
    Lowercase, strip filler words/punctuation, condense spaces.
    Examples:
      'Close the YouTube app.' -> 'youtube'
      'exit   chrome!!!'       -> 'chrome'
    """
    s = s.lower()
    for w in ["the", "a", "an", "app", "application", "please"]:
        s = s.replace(f" {w} ", " ")
        s = s.replace(f" {w}", " ")
        s = s.replace(f"{w} ", " ")
    s = re.sub(r"[^a-z0-9+.\s]", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s

def _is_running(proc_name: str) -> bool:
    try:
        out = subprocess.check_output(["tasklist"], creationflags=subprocess.CREATE_NO_WINDOW).decode(errors="ignore")
        return proc_name.lower() in out.lower()
    except Exception:
        return False

def close_application(app_name):
    try:
        app_name = _normalize_app_name(app_name)

        # Map common phrases to actual Windows process names
        candidates = {
            # Browsers (YouTube/Hotstar are tabs -> kill the browser)
            "youtube": ["msedge.exe", "chrome.exe", "firefox.exe"],
            "yt": ["msedge.exe", "chrome.exe", "firefox.exe"],
            "hotstar": ["msedge.exe", "chrome.exe", "firefox.exe"],
            "jiohotstar": ["msedge.exe", "chrome.exe", "firefox.exe"],

            "edge": ["msedge.exe"],
            "microsoft edge": ["msedge.exe"],
            "chrome": ["chrome.exe"],
            "google chrome": ["chrome.exe"],
            "firefox": ["firefox.exe"],

            "word": ["WINWORD.EXE"],
            "excel": ["EXCEL.EXE"],
            "powerpoint": ["POWERPNT.EXE"],
            "ppt": ["POWERPNT.EXE"],

            "notepad": ["notepad.exe"],
            "notepad++": ["notepad++.exe"],
            "vlc": ["vlc.exe"],
            "spotify": ["spotify.exe"],
            "zoom": ["Zoom.exe"],
            "teams": ["Teams.exe"],
            "whatsapp": ["WhatsApp.exe"],
        }

        # Fuzzy map: pick the first key contained in the request
        proc_list = None
        for key, procs in candidates.items():
            if key in app_name:
                proc_list = procs
                break

        # If no fuzzy key matched, try a best-guess from the cleaned name
        if not proc_list:
            guess = app_name.replace(" ", "") + ".exe"
            proc_list = [guess]

        killed_any = False
        errs = []

        for proc in proc_list:
            if _is_running(proc):
                try:
                    # Kill the whole process tree
                    subprocess.run(
                        ["taskkill", "/F", "/T", "/IM", proc],
                        check=True,
                        creationflags=subprocess.CREATE_NO_WINDOW
                    )
                    killed_any = True
                except subprocess.CalledProcessError as e:
                    errs.append(f"{proc}: {e}")

        if killed_any:
            return f"✅ Closed: {', '.join(proc_list)}"
        else:
            if errs:
                return "⚠️ Tried to close but failed: " + "; ".join(errs)
            return f"⚠️ Couldn’t find a running process to close for “{app_name}”. If this is a website tab, close the browser (e.g., say: ‘close edge’)."

    except Exception as e:
        return f"⚠️ Close error: {e}"

# ------------------ WhatsApp Automation ------------------
def open_whatsapp_and_send(contact_name, message="", call_type=None):
    try:
        subprocess.Popen(["start", "whatsapp:"], shell=True)
    except:
        return "⚠️ Could not open WhatsApp"

    time.sleep(8)
    pyautogui.hotkey('alt', 'tab')
    time.sleep(1)
    pyautogui.hotkey('ctrl', 'f')
    pyautogui.write(contact_name)
    pyautogui.press('enter')
    time.sleep(2)

    if message:
        pyautogui.write(message)
        pyautogui.press('enter')
    return f"✅ WhatsApp action completed for {contact_name}"

# ------------------ Folder Creation ------------------
def create_folder(folder_name="New Folder"):
    desktop_path = os.path.join(os.path.expanduser("~"), "Desktop")
    path = os.path.join(desktop_path, folder_name)
    try:
        os.makedirs(path, exist_ok=True)
        return f"✅ Folder '{folder_name}' created on Desktop at {path}"
    except Exception as e:
        return f"⚠️ Failed to create folder: {e}"

# ------------------ PDF from Recent Downloads ------------------
def create_pdf_from_recent_downloads(minutes=5, folder_path=None):
    try:
        if not folder_path:
            folder_path = os.path.join(os.path.expanduser("~"), "Downloads")

        now = time.time()
        cutoff = now - (minutes * 60)

        images = [f for f in os.listdir(folder_path)
                  if f.lower().endswith(('.png', '.jpg', '.jpeg'))
                  and os.path.getctime(os.path.join(folder_path, f)) >= cutoff]

        if not images:
            return f"⚠️ No images found in Downloads from the last {minutes} minutes."

        images.sort(key=lambda x: os.path.getctime(os.path.join(folder_path, x)))

        pil_images = []
        for img in images:
            im = Image.open(os.path.join(folder_path, img))
            if im.mode == "RGBA":
                im = im.convert("RGB")
            pil_images.append(im)

        pdf_path = os.path.join(folder_path, f"RecentDownloadsPDF_{int(time.time())}.pdf")
        pil_images[0].save(pdf_path, save_all=True, append_images=pil_images[1:])
        os.startfile(pdf_path)
        return f"✅ PDF created from {len(images)} recent images: {pdf_path}"

    except Exception as e:
        return f"⚠️ Failed to create PDF: {str(e)}"

# ------------------ Command Handling ------------------
def automation_commands(cmd):
    cmd = cmd.lower()

    # Speak the command acknowledgement
    speak(f"Okay, I will {cmd}")

    if "create a folder" in cmd:
        folder_name = cmd.replace("create a folder", "").strip()
        return create_folder(folder_name if folder_name else "New Folder")

    if "create pdf from recent downloads" in cmd:
        return create_pdf_from_recent_downloads()

    # Close apps  ✅ normalized + resilient
    if "close" in cmd or "exit" in cmd or "shut" in cmd:
        app_name = cmd
        for v in ["close", "exit", "shut", "turn off", "kill", "stop"]:
            app_name = app_name.replace(v, " ")
        app_name = _normalize_app_name(app_name)
        return close_application(app_name if app_name else "")

    # AI content
    for key in ["application","essay","letter","story","report","speech","email"]:
        if key in cmd:
            topic = cmd.replace(f"write {key}", "").strip()
            content = write_content(f"Write a {key} about {topic}")
            open_notepad_with_content(content)
            return f"✅ {key.capitalize()} opened in Notepad."

    # PPT
    if "ppt" in cmd or "presentation" in cmd:
        topic = cmd.replace("ppt", "").replace("presentation", "").strip()
        return create_ppt(topic if topic else "Topic")

    # WhatsApp
    if "whatsapp" in cmd:
        message = ""
        contact_name = cmd.replace("whatsapp", "").replace("send message to", "").strip()
        if "message" in contact_name:
            parts = contact_name.split("message", 1)
            contact_name = parts[0].strip()
            message = parts[1].strip()
        return open_whatsapp_and_send(contact_name, message)

    # Open apps
    if any(x in cmd for x in ["open", "launch"]):
        app_name = cmd.replace("open", "").replace("launch", "").strip()
        return open_application(app_name if app_name else "")

    return "⚠️ Command not recognized."

# ------------------ Run ------------------
if __name__ == "__main__":
    speak("Automation module is ready. Say 'exit' to quit.")
    print("🎛 Automation Ready. Listening for voice commands...")

    while True:
        cmd = listen_command()
        if cmd in ["exit","quit","bye"]:
            speak("Exiting Automation module. Goodbye!")
            break
        if cmd:
            output = automation_commands(cmd)
            print(f"➡️ {output}")
            speak(output)
